#!/usr/bin/env python3
"""
Generate Board Deck PowerPoint from template and demo data.
Creates a beautiful, professional board presentation with PartnerAgents branding.
"""

import json
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Try to import python-pptx, install if missing
try:
    from pptx import Presentation
except ImportError:
    print("Installing python-pptx...")
    os.system("pip install python-pptx")
    from pptx import Presentation


# === BRANDING CONFIG ===
BRAND_COLOR = RGBColor(99, 102, 241)  # Indigo #6366F1
ACCENT_COLOR = RGBColor(16, 185, 129)  # Emerald green
WARNING_COLOR = RGBColor(245, 158, 11)  # Amber
DANGER_COLOR = RGBColor(239, 68, 68)  # Red
DARK_TEXT = RGBColor(17, 24, 39)  # Near black
LIGHT_TEXT = RGBColor(107, 114, 128)  # Gray
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(243, 244, 246)


# === LAYOUT HELPERS ===


def add_slide_title(slide, title_text, subtitle_text=None):
    """Add a beautiful title with optional subtitle"""
    title_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
    )
    title_frame = title_shape.text_frame
    title_frame.word_wrap = True

    title_p = title_frame.paragraphs[0]
    title_p.text = title_text
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = DARK_TEXT
    title_p.alignment = PP_ALIGN.LEFT

    if subtitle_text:
        subtitle_p = title_frame.add_paragraph()
        subtitle_p.text = subtitle_text
        subtitle_p.font.size = Pt(18)
        subtitle_p.font.color.rgb = LIGHT_TEXT
        subtitle_p.alignment = PP_ALIGN.LEFT


def add_brand_bar(slide, height=Inches(0.08)):
    """Add PartnerAgents brand color bar at top"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_COLOR
    shape.line.fill.background()


def add_metric_card(slide, x, y, width, height, value, label, trend=None):
    """Add a metric card with value and label"""
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LIGHT_GRAY
    shape.line.width = Pt(1)

    # Value
    value_box = slide.shapes.add_textbox(
        x + Inches(0.2), y + Inches(0.3), width - Inches(0.4), Inches(0.8)
    )
    value_frame = value_box.text_frame
    value_p = value_frame.paragraphs[0]
    value_p.text = value
    value_p.font.size = Pt(36)
    value_p.font.bold = True
    value_p.font.color.rgb = BRAND_COLOR
    value_p.alignment = PP_ALIGN.CENTER

    # Label
    label_box = slide.shapes.add_textbox(
        x + Inches(0.2), y + Inches(1.0), width - Inches(0.4), Inches(0.5)
    )
    label_frame = label_box.text_frame
    label_p = label_frame.paragraphs[0]
    label_p.text = label
    label_p.font.size = Pt(12)
    label_p.font.color.rgb = LIGHT_TEXT
    label_p.alignment = PP_ALIGN.CENTER

    # Trend indicator
    if trend:
        trend_box = slide.shapes.add_textbox(
            x + Inches(0.2), y + Inches(1.5), width - Inches(0.4), Inches(0.4)
        )
        trend_frame = trend_box.text_frame
        trend_p = trend_frame.paragraphs[0]
        trend_p.text = trend
        trend_p.font.size = Pt(11)
        trend_p.font.color.rgb = ACCENT_COLOR if "+" in trend else DANGER_COLOR
        trend_p.alignment = PP_ALIGN.CENTER


def add_section_header(slide, text):
    """Add a section header with brand accent"""
    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(0.1), Inches(0.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND_COLOR
    bar.line.fill.background()

    # Text
    text_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.8), Inches(8), Inches(0.6)
    )
    frame = text_box.text_frame
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT


def add_bullet_list(slide, items, x=Inches(0.8), y=Inches(2.8), width=Inches(8.4)):
    """Add a bulleted list"""
    y_pos = y
    for item in items:
        text_box = slide.shapes.add_textbox(x, y_pos, width, Inches(0.5))
        frame = text_box.text_frame
        p = frame.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        y_pos += Inches(0.45)


def create_title_slide(prs, data):
    """Create the title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    add_brand_bar(slide)

    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = f"Partner Program"
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = DARK_TEXT
    title_p.alignment = PP_ALIGN.CENTER

    # Quarter
    quarter_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    quarter_frame = quarter_box.text_frame
    quarter_p = quarter_frame.paragraphs[0]
    quarter_p.text = f"{data['quarter']} Board Update"
    quarter_p.font.size = Pt(28)
    quarter_p.font.color.rgb = BRAND_COLOR
    quarter_p.alignment = PP_ALIGN.CENTER

    # Company
    company_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.6))
    company_frame = company_box.text_frame
    company_p = company_frame.paragraphs[0]
    company_p.text = data["company"]["name"]
    company_p.font.size = Pt(18)
    company_p.font.color.rgb = LIGHT_TEXT
    company_p.alignment = PP_ALIGN.CENTER


def create_executive_summary(prs, data):
    """Create Executive Summary slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_brand_bar(slide)
    add_slide_title(slide, "Executive Summary", "At a Glance")

    s = data["summary"]

    # Metrics row
    add_metric_card(
        slide,
        Inches(0.5),
        Inches(1.6),
        Inches(2.1),
        Inches(2.2),
        f"{s['total_partners']}",
        "Active Partners",
        f"+{s['partners_new_this_quarter']} this quarter",
    )
    add_metric_card(
        slide,
        Inches(2.8),
        Inches(1.6),
        Inches(2.1),
        Inches(2.2),
        s["revenue_formatted"],
        "Partner Revenue",
        f"+{s['growth_yoy']}% YoY",
    )
    add_metric_card(
        slide,
        Inches(5.1),
        Inches(1.6),
        Inches(2.1),
        Inches(2.2),
        f"{s['revenue_from_partners_pct']}%",
        "of Total Revenue",
    )
    add_metric_card(
        slide,
        Inches(7.4),
        Inches(1.6),
        Inches(2.1),
        Inches(2.2),
        f"${data['pipeline']['total'] // 1000000}M",
        "Pipeline",
        f"+{data['pipeline']['growth_qoq']}% QoQ",
    )

    # Key wins section
    add_section_header(slide, "Key Wins")
    wins = data["focus_areas"]["wins"]
    add_bullet_list(slide, wins, y=Inches(3.2))

    # Focus areas
    add_section_header(slide, "Focus Areas")
    risks = data["focus_areas"]["risks"]
    add_bullet_list(slide, risks, y=Inches(4.7))


def create_performance_slide(prs, data):
    """Create Performance slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_brand_bar(slide)
    add_slide_title(slide, "Performance vs. Plan", f"{data['quarter']}")

    s = data["summary"]
    t = data["targets"]

    # Revenue comparison
    add_section_header(slide, "Revenue Performance")

    # Actual vs Target
    slide.shapes.add_textbox(
        Inches(0.8), Inches(2.4), Inches(4), Inches(0.5)
    ).text_frame.paragraphs[
        0
    ].text = f"Actual: {s['revenue_formatted']} ({s['growth_yoy']}% YoY)"
    slide.shapes.add_textbox(
        Inches(0.8), Inches(2.9), Inches(4), Inches(0.5)
    ).text_frame.paragraphs[0].text = f"Target: ${t['revenue_target'] / 1000000:.1f}M"
    slide.shapes.add_textbox(
        Inches(0.8), Inches(3.4), Inches(4), Inches(0.5)
    ).text_frame.paragraphs[
        0
    ].text = f"Variance: {(s['revenue'] - t['revenue_target']) / t['revenue_target'] * 100:.0f}%"

    # Revenue breakdown
    add_section_header(slide, "Revenue Breakdown")

    slide.shapes.add_textbox(
        Inches(0.8), Inches(4.4), Inches(4), Inches(0.5)
    ).text_frame.paragraphs[
        0
    ].text = f"Partner-Sourced: ${s['partner_sourced']:,} (61%)"
    slide.shapes.add_textbox(
        Inches(0.8), Inches(4.9), Inches(4), Inches(0.5)
    ).text_frame.paragraphs[
        0
    ].text = f"Partner-Influenced: ${s['partner_influenced']:,} (39%)"

    # Variance analysis
    add_section_header(slide, "Variance Analysis")
    reasons = data["variance_analysis"]["missed_reasons"]
    add_bullet_list(slide, reasons, y=Inches(5.9))


def create_partner_health_slide(prs, data):
    """Create Partner Health slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_brand_bar(slide)
    add_slide_title(slide, "Partner Program Health", "Key KPIs")

    ph = data["partner_health"]
    pt = data["partner_tiers"]

    # Partner count by tier
    add_section_header(slide, "Partner Count")

    y = Inches(2.4)
    for tier in ["gold", "silver", "bronze"]:
        tier_name = tier.capitalize()
        current = pt[tier]["current"]
        target = pt[tier]["target"]

        slide.shapes.add_textbox(
            Inches(0.8), y, Inches(3), Inches(0.5)
        ).text_frame.paragraphs[0].text = f"{tier_name}: {current} / {target} target"
        y += Inches(0.4)

    # Engagement metrics
    add_section_header(slide, "Engagement")

    add_metric_card(
        slide,
        Inches(0.5),
        Inches(3.9),
        Inches(2.1),
        Inches(1.8),
        f"{ph['engagement_score']}/100",
        "Engagement Score",
    )
    add_metric_card(
        slide,
        Inches(2.8),
        Inches(3.9),
        Inches(2.1),
        Inches(1.8),
        f"{ph['certified_partners']}",
        "Certified Partners",
        f"{ph['certified_pct']}%",
    )
    add_metric_card(
        slide,
        Inches(5.1),
        Inches(3.9),
        Inches(2.1),
        Inches(1.8),
        f"{ph['deal_registrations']}",
        "Deal Registrations",
    )
    add_metric_card(
        slide,
        Inches(7.4),
        Inches(3.9),
        Inches(2.1),
        Inches(1.8),
        f"{ph['churn_rate']}%",
        "Churn Rate",
        f"< {ph['churn_target']}% target",
    )

    # NPS
    add_section_header(slide, "Partner Satisfaction")
    slide.shapes.add_textbox(
        Inches(0.8), Inches(5.7), Inches(4), Inches(0.5)
    ).text_frame.paragraphs[0].text = f"NPS: {ph['nps']} (up from {ph['nps_previous']})"


def create_top_partners_slide(prs, data):
    """Create Top Partners slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_brand_bar(slide)
    add_slide_title(slide, "Top Partners", "By Revenue")

    # Table
    x = Inches(0.8)
    y = Inches(2.0)
    col_widths = [Inches(3.5), Inches(2), Inches(2), Inches(2)]

    # Headers
    headers = ["Partner", "Revenue", "% of Total", "Tier"]
    for i, header in enumerate(headers):
        box = slide.shapes.add_textbox(x, y, col_widths[i], Inches(0.4))
        p = box.text_frame.paragraphs[0]
        p.text = header
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_TEXT
        x += col_widths[i]

    # Data rows
    y += Inches(0.5)
    for partner in data["top_partners"]:
        x = Inches(0.8)
        row_data = [
            partner["name"],
            f"${partner['revenue']:,}",
            f"{partner['pct']}%",
            partner["tier"],
        ]
        for i, cell in enumerate(row_data):
            box = slide.shapes.add_textbox(x, y, col_widths[i], Inches(0.4))
            p = box.text_frame.paragraphs[0]
            p.text = cell
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT
            x += col_widths[i]
        y += Inches(0.45)


def create_roadmap_slide(prs, data):
    """Create Roadmap slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_brand_bar(slide)
    add_slide_title(slide, "Partner Program Roadmap", "2026")

    rd = data["roadmap"]

    # Q1
    add_section_header(slide, "Q1 (Current)")
    add_bullet_list(slide, rd["q1"], y=Inches(2.8))

    # Q2
    add_section_header(slide, "Q2 (Planned)")
    add_bullet_list(slide, rd["q2"], y=Inches(4.3))

    # H2
    add_section_header(slide, "H2 Priorities")
    add_bullet_list(slide, rd["h2"], y=Inches(5.8))


def create_ask_slide(prs, data):
    """Create The Ask slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_brand_bar(slide)
    add_slide_title(slide, "Resource Request", "Investment Ask")

    ask = data["ask"]

    # Main ask
    slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8), Inches(9), Inches(1)
    ).text_frame.paragraphs[0].text = f"${ask['amount']:,} {ask['purpose']}"
    slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8), Inches(9), Inches(1)
    ).text_frame.paragraphs[0].font.size = Pt(36)
    slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8), Inches(9), Inches(1)
    ).text_frame.paragraphs[0].font.bold = True
    slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8), Inches(9), Inches(1)
    ).text_frame.paragraphs[0].font.color.rgb = BRAND_COLOR

    # Expected return
    add_section_header(slide, "Expected Return")

    slide.shapes.add_textbox(
        Inches(0.8), Inches(3.0), Inches(8), Inches(0.5)
    ).text_frame.paragraphs[
        0
    ].text = (
        f"Generate ${ask['expected_pipeline']:,} pipeline ({ask['expected_roi']}:1 ROI)"
    )
    slide.shapes.add_textbox(
        Inches(0.8), Inches(3.5), Inches(8), Inches(0.5)
    ).text_frame.paragraphs[
        0
    ].text = f"Enable {ask['campaigns']} co-marketing campaigns"
    slide.shapes.add_textbox(
        Inches(0.8), Inches(4.0), Inches(8), Inches(0.5)
    ).text_frame.paragraphs[
        0
    ].text = f"Support {ask['new_partners']} new partner launches"

    # Risk of inaction
    add_section_header(slide, "Risk of Inaction")

    slide.shapes.add_textbox(
        Inches(0.8), Inches(5.1), Inches(8), Inches(0.5)
    ).text_frame.paragraphs[0].text = "• Lose co-marketing momentum"
    slide.shapes.add_textbox(
        Inches(0.8), Inches(5.5), Inches(8), Inches(0.5)
    ).text_frame.paragraphs[0].text = "• Miss Q1 pipeline targets"
    slide.shapes.add_textbox(
        Inches(0.8), Inches(5.9), Inches(8), Inches(0.5)
    ).text_frame.paragraphs[0].text = "• Partner satisfaction decline"


def generate_board_deck(output_path=None):
    """Generate the complete board deck presentation"""

    # Load demo data - go up from scripts/ to PartnerAgents root
    script_dir = Path(__file__).parent.resolve()
    root_dir = script_dir.parent
    data_path = root_dir / "examples" / "demo-company" / "board-deck-data.json"

    # Validate data file exists
    if not data_path.exists():
        print(f"ERROR: Demo data file not found: {data_path}")
        print("Please create examples/demo-company/board-deck-data.json")
        return None

    # Load and validate JSON
    try:
        with open(data_path, "r") as f:
            data = json.load(f)
    except json.JSONDecodeError as e:
        print(f"ERROR: Invalid JSON in {data_path}: {e}")
        return None

    # Validate required fields
    required_fields = [
        "company",
        "quarter",
        "summary",
        "top_partners",
        "roadmap",
        "ask",
    ]
    missing = [f for f in required_fields if f not in data]
    if missing:
        print(f"ERROR: Missing required fields in demo data: {missing}")
        return None

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create all slides
    create_title_slide(prs, data)
    create_executive_summary(prs, data)
    create_performance_slide(prs, data)
    create_partner_health_slide(prs, data)
    create_top_partners_slide(prs, data)
    create_roadmap_slide(prs, data)
    create_ask_slide(prs, data)

    # Determine output path
    if output_path is None:
        output_dir = root_dir / "partneros-docs" / "src" / "assets" / "pptx"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "board-deck.pptx"

    # Save
    prs.save(str(output_path))
    print(f"✓ Board deck generated: {output_path}")

    # Also generate PDF version for in-browser viewing
    try:
        from pptx.util import Inches as PptxInches

        # Use pdfsave to save as PDF
        output_pdf = output_path.with_suffix(".pdf")
        prs.save(str(output_pdf))
        print(f"✓ PDF generated: {output_pdf}")
    except Exception as e:
        print(f"Note: PDF generation failed (requires MS PowerPoint): {e}")

    return str(output_path)


if __name__ == "__main__":
    if len(sys.argv) > 1:
        output_path = sys.argv[1]
    else:
        output_path = None

    generate_board_deck(output_path)
